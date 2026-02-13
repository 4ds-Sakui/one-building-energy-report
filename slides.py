#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
スライド作成モジュール (Streamlit対応版 v1.3)
モデル建物法と標準入力法の自動切り替え対応
標準入力法の「ちら見せ」と組織自立診断への誘導を実装
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import io
import tempfile
import os

# report_generator.pyからカラー定義をインポート
from report_generator import (
    COLOR_MAIN, COLOR_RED, COLOR_GREEN,
    get_bei_label, get_bpi_label,
    generate_improvement_roadmap
)

# 新しいカラー定義 (HTMLスライドと同期)
COLOR_ACCENT = RGBColor(244, 162, 97) # #F4A261
COLOR_LIGHT_BLUE = RGBColor(231, 243, 255) # #e7f3ff
COLOR_LIGHT_YELLOW = RGBColor(255, 255, 243) # #fffff3
COLOR_LIGHT_RED = RGBColor(255, 243, 243) # #fff3f3
COLOR_LIGHT_GREEN = RGBColor(243, 255, 243) # #f3fff3
COLOR_GRAY = RGBColor(153, 153, 153) # #999999
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)

def create_presentation(data, chart_stacked_bytes, chart_pie_bytes, chart_bei_bytes):
    """
    PowerPointプレゼンテーションを作成してBytesIOで返す
    """
    prs = Presentation()
    
    # 16:9ワイドスクリーン
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    add_title_slide_tech_report_style(prs, data)
    add_summary_slide(prs, data)
    
    calc_method = data.get("calculation_method", "standard_input")
    is_model = (calc_method == "model_building")
    
    # 標準入力法の場合のみ、詳細グラフを追加
    if not is_model:
        # 一時ファイルにグラフを保存してスライドに追加
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_stacked:
            tmp_stacked.write(chart_stacked_bytes.read())
            tmp_stacked_path = tmp_stacked.name
        chart_stacked_bytes.seek(0)
        
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_pie:
            tmp_pie.write(chart_pie_bytes.read())
            tmp_pie_path = tmp_pie.name
        chart_pie_bytes.seek(0)
        
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(slide3, "エネルギー消費性能の詳細分析")
        slide3.shapes.add_picture(tmp_stacked_path, Inches(0.3), Inches(0.95), width=Inches(9.4))
        
        slide4 = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(slide4, "設備別一次エネルギー消費量の比較")
        slide4.shapes.add_picture(tmp_pie_path, Inches(0.3), Inches(1.1), width=Inches(9.4))
        
        os.unlink(tmp_stacked_path)
        os.unlink(tmp_pie_path)
    
    add_envelope_worst_analysis_slide(prs, data)
    
    # BEI比較グラフを一時ファイルに保存
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_bei:
        tmp_bei.write(chart_bei_bytes.read())
        tmp_bei_path = tmp_bei.name
    chart_bei_bytes.seek(0)
    
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    bei_label = get_bei_label(calc_method)
    add_slide_title(slide6, f"用途別エネルギー消費傾向: {bei_label}分析")
    slide6.shapes.add_picture(tmp_bei_path, Inches(0.3), Inches(1.05), width=Inches(9.4))
    
    os.unlink(tmp_bei_path)
    
    # モデル建物法の場合に標準入力法への誘導スライドを追加
    if is_model:
        add_standard_input_teaser_slide(prs, data)

    add_improvement_roadmap_slide(prs, data)
    add_organizational_diagnosis_slide(prs, data)
    
    # BytesIOに保存
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes

def add_slide_title(slide, title_text):
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(22)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_MAIN
    title_para.font.name = 'Noto Sans JP'

def add_title_slide_tech_report_style(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    center_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(0.8), Inches(5), Inches(4))
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = COLOR_MAIN
    center_box.line.fill.background()
    
    logo_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.3), Inches(3), Inches(0.4))
    logo_box.text_frame.text = "one building"
    logo_box.text_frame.paragraphs[0].font.size = Pt(24)
    logo_box.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    logo_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    building_box = slide.shapes.add_textbox(Inches(3), Inches(2.0), Inches(4), Inches(0.4))
    building_box.text_frame.text = data['building_name']
    building_box.text_frame.paragraphs[0].font.size = Pt(22)
    building_box.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    building_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(3), Inches(2.6), Inches(4), Inches(0.5))
    title_box.text_frame.text = "技術レポート"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    date_box = slide.shapes.add_textbox(Inches(3.5), Inches(3.5), Inches(3), Inches(0.35))
    date_box.text_frame.text = datetime.now().strftime('%Y.%m.%d')
    date_box.text_frame.paragraphs[0].font.size = Pt(16)
    date_box.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    date_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "1. 総合評価サマリー: 現状と経営リスク")
    
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
    info_box.text_frame.text = f"建物名称: {data['building_name']}  |  所在地: {data['location']}  |  延べ面積: {data['total_area']} m²"
    info_box.text_frame.paragraphs[0].font.size = Pt(11)
    info_box.text_frame.paragraphs[0].font.color.rgb = COLOR_GRAY
    
    is_compliant = (data['bei_total'] <= 1.0)
    status_color = COLOR_GREEN if is_compliant else COLOR_RED
    status_bg_color = COLOR_LIGHT_GREEN if is_compliant else COLOR_LIGHT_RED
    status_text = f"診断結果: {('基準適合' if is_compliant else '基準非適合')}"
    
    # 総合判定結果
    res_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(0.5))
    res_box.text_frame.text = status_text
    res_box.text_frame.paragraphs[0].font.size = Pt(18)
    res_box.text_frame.paragraphs[0].font.bold = True
    res_box.text_frame.paragraphs[0].font.color.rgb = status_color
    
    # 経営リスク評価
    risk_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(9), Inches(3.0))
    risk_shape.fill.solid()
    risk_shape.fill.fore_color.rgb = status_bg_color
    risk_shape.line.color.rgb = status_color
    
    tf = risk_shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p1 = tf.paragraphs[0]
    p1.text = f"▲{('優位性' if is_compliant else '重要')}: 経営影響の特定"
    p1.font.bold = True
    p1.font.color.rgb = status_color
    p1.font.size = Pt(14)
    
    # 法的リスク
    p2 = tf.add_paragraph()
    p2.text = "● 法的リスク: "
    if is_compliant:
        p2.text += "基準適合。建築確認申請がスムーズに進められます。"
    else:
        p2.text += "基準非適合。改正省エネ法に基づき建築確認が受理されない恐れがあります。"
    p2.font.size = Pt(12)
    p2.font.name = 'Noto Sans JP'

    # 事業リスク
    p3 = tf.add_paragraph()
    p3.text = "● 事業リスク: "
    if is_compliant:
        p3.text += "光熱費削減による運用コストの低減、企業イメージ向上。"
    else:
        p3.text += "高い光熱費による運用コストの逼迫、競争力低下。"
    p3.font.size = Pt(12)
    p3.font.name = 'Noto Sans JP'

    # 資産価値リスク
    p4 = tf.add_paragraph()
    p4.text = "● 資産価値リスク: "
    if is_compliant:
        p4.text += "ZEB認定取得の可能性、ESG投資基準への適合、不動産価値向上。"
    else:
        p4.text += "ZEB化遅延による不動産価値の低下、市場評価の悪化。"
    p4.font.size = Pt(12)
    p4.font.name = 'Noto Sans JP'

def add_envelope_worst_analysis_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    calc_method = data.get('calculation_method', 'standard_input')
    bpi_label = get_bpi_label(calc_method)
    add_slide_title(slide, f"2. 外皮性能評価 ({bpi_label}) と改善ポイント")
    
    bpi_val = data.get('bpi', 1.0)
    pal_design = data.get('pal_design', 1.0)
    pal_standard = data.get('pal_standard', 1.0)

    # BPI値の表示
    bpi_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(1.0))
    bpi_tf = bpi_box.text_frame
    bpi_tf.text = f"{bpi_label} 値: {bpi_val:.2f}"
    bpi_tf.paragraphs[0].font.size = Pt(36)
    bpi_tf.paragraphs[0].font.bold = True
    bpi_tf.paragraphs[0].font.color.rgb = COLOR_MAIN

    # BPI評価コメント
    bpi_comment_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.5))
    bpi_comment_tf = bpi_comment_box.text_frame
    bpi_comment_tf.text = f"外皮性能は{bpi_val:.2f}で、基準値1.0{('以下' if bpi_val <= 1.0 else '超過')}です。"
    bpi_comment_tf.paragraphs[0].font.size = Pt(14)
    bpi_comment_tf.paragraphs[0].font.color.rgb = COLOR_GREEN if bpi_val <= 1.0 else COLOR_RED

    # 改善のポイント
    points_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(2.5))
    points_tf = points_box.text_frame
    points_tf.text = "改善のポイント:"
    points_tf.paragraphs[0].font.bold = True
    points_tf.paragraphs[0].font.size = Pt(16)

    p_window = points_tf.add_paragraph()
    p_window.text = "● 窓のLow-E化（日射熱取得率の低減）"
    p_window.font.size = Pt(12)

    p_insulation = points_tf.add_paragraph()
    p_insulation.text = "● 断熱材の厚さ・性能向上"
    p_insulation.font.size = Pt(12)

    p_thermal_bridge = points_tf.add_paragraph()
    p_thermal_bridge.text = "● 熱橋部分の対策"
    p_thermal_bridge.font.size = Pt(12)

    p_airtightness = points_tf.add_paragraph()
    p_airtightness.text = "● 気密性の向上"
    p_airtightness.font.size = Pt(12)

    # 標準入力法の場合のみワースト室分析を追加
    if not calc_method == 'model_building' and data.get('worst_rooms'):
        worst_rooms_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.0), Inches(4.0), Inches(4.0))
        worst_rooms_tf = worst_rooms_box.text_frame
        worst_rooms_tf.text = "ワースト要因の詳細分析:"
        worst_rooms_tf.paragraphs[0].font.bold = True
        worst_rooms_tf.paragraphs[0].font.size = Pt(16)
        worst_rooms_tf.paragraphs[0].font.color.rgb = COLOR_RED

        for i, room in enumerate(data.get('worst_rooms', [])[:3]): # 最大3つまで表示
            p_room = worst_rooms_tf.add_paragraph()
            p_room.text = f"・{room.get('name', '不明')}: {room.get('factor', '不明')} → {room.get('improvement', '改善策検討中')}"
            p_room.font.size = Pt(10)
            p_room.font.color.rgb = COLOR_BLACK

def add_standard_input_teaser_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "標準入力法による詳細分析のご紹介")

    # 導入テキスト
    intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.8))
    intro_tf = intro_box.text_frame
    intro_tf.text = "モデル建物法では得られない、さらに詳細な省エネ分析で、貴社のZEB化を強力にサポートします。"
    intro_tf.paragraphs[0].font.size = Pt(16)
    intro_tf.paragraphs[0].font.bold = True
    intro_tf.paragraphs[0].font.color.rgb = COLOR_MAIN

    # 詳細分析のメリットをグリッド形式で表示
    # 1列目
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0), Inches(4.0), Inches(1.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    box1.line.fill.background()
    tf1 = box1.text_frame
    tf1.text = "🏢 室別負荷分析\n各室の熱負荷を特定し、改修優先度を判定します。"
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.size = Pt(14)
    tf1.paragraphs[1].font.size = Pt(10)

    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.7), Inches(4.0), Inches(1.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    box2.line.fill.background()
    tf2 = box2.text_frame
    tf2.text = "📊 エネルギー消費分布\n設備別の消費量を詳細に把握し、最適な改善策を提案します。"
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[1].font.size = Pt(10)

    # 2列目
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(2.0), Inches(4.0), Inches(1.5))
    box3.fill.solid()
    box3.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    box3.line.fill.background()
    tf3 = box3.text_frame
    tf3.text = "💰 LCC分析\n初期投資と運用コストを総合的に評価します。"
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[1].font.size = Pt(10)

    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(3.7), Inches(4.0), Inches(1.5))
    box4.fill.solid()
    box4.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    box4.line.fill.background()
    tf4 = box4.text_frame
    tf4.text = "🎯 ZEB化ロードマップ\n段階的な改修計画と投資効果を可視化します。"
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].font.size = Pt(14)
    tf4.paragraphs[1].font.size = Pt(10)

    # 誘導テキスト
    guidance_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(0.3))
    guidance_tf = guidance_box.text_frame
    guidance_tf.text = "これらの詳細分析は、貴社のZEB化戦略を加速させ、持続可能な経営を実現するための強力なツールとなります。"
    guidance_tf.paragraphs[0].font.size = Pt(12)
    guidance_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    guidance_tf.paragraphs[0].font.color.rgb = COLOR_BLACK

def add_improvement_roadmap_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "4. ZEB化への改善ロードマップ")
    
    roadmap = generate_improvement_roadmap(data)
    
    # ロードマップの各ステップを配置
    for i, step in enumerate(roadmap):
        # カードの背景
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*2.3), Inches(1.2), Inches(2.1), Inches(4.0))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_WHITE
        box.line.color.rgb = COLOR_MAIN
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        # STEP番号
        p_step = tf.paragraphs[0]
        p_step.text = step['step']
        p_step.font.bold = True
        p_step.font.size = Pt(24)
        p_step.font.color.rgb = COLOR_MAIN
        p_step.alignment = PP_ALIGN.CENTER

        # タイトル
        p_title = tf.add_paragraph()
        p_title.text = step['title']
        p_title.font.bold = True
        p_title.font.size = Pt(14)
        p_title.font.color.rgb = COLOR_BLACK
        p_title.alignment = PP_ALIGN.CENTER

        # 説明
        p_desc = tf.add_paragraph()
        p_desc.text = step['desc']
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = COLOR_BLACK
        p_desc.alignment = PP_ALIGN.LEFT

    # 追加のメッセージ
    message_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(0.3))
    message_tf = message_box.text_frame
    message_tf.text = "これらのステップを通じて、貴社のZEB化を段階的に実現し、持続可能な社会に貢献します。"
    message_tf.paragraphs[0].font.size = Pt(12)
    message_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    message_tf.paragraphs[0].font.color.rgb = COLOR_BLACK

def add_organizational_diagnosis_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "組織自立診断へのご招待")

    # 導入テキスト
    intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.8))
    intro_tf = intro_box.text_frame
    intro_tf.text = "貴社の省エネ・ZEB化推進における組織的な課題と機会を特定し、最適な戦略をご提案します。"
    intro_tf.paragraphs[0].font.size = Pt(16)
    intro_tf.paragraphs[0].font.bold = True
    intro_tf.paragraphs[0].font.color.rgb = COLOR_MAIN

    # 診断項目
    diagnosis_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.5))
    diagnosis_box.fill.solid()
    diagnosis_box.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    diagnosis_box.line.fill.background()
    diagnosis_tf = diagnosis_box.text_frame
    diagnosis_tf.text = "診断項目（約5分）:"
    diagnosis_tf.paragraphs[0].font.bold = True
    diagnosis_tf.paragraphs[0].font.size = Pt(14)

    p_commit = diagnosis_tf.add_paragraph()
    p_commit.text = "・経営層のコミットメント"
    p_commit.font.size = Pt(12)

    p_coop = diagnosis_tf.add_paragraph()
    p_coop.text = "・技術部門と営業部門の連携"
    p_coop.font.size = Pt(12)

    p_team = diagnosis_tf.add_paragraph()
    p_team.text = "・推進チームの有無"
    p_team.font.size = Pt(12)

    p_seminar = diagnosis_tf.add_paragraph()
    p_seminar.text = "・過去のセミナー受講歴"
    p_seminar.font.size = Pt(12)

    # 誘導テキスト
    guidance_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.8))
    guidance_tf = guidance_box.text_frame
    guidance_tf.text = "診断結果に基づき、貴社に最適なソリューションと具体的なアクションプランをご提案いたします。\nぜひ、この機会に貴社の組織力を診断し、ZEB化への道を加速させましょう。"
    guidance_tf.paragraphs[0].font.size = Pt(12)
    guidance_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    guidance_tf.paragraphs[0].font.color.rgb = COLOR_BLACK
    guidance_tf.paragraphs[1].font.size = Pt(12)
    guidance_tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    guidance_tf.paragraphs[1].font.color.rgb = COLOR_BLACK
